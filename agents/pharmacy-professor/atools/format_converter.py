#!/usr/bin/env python3
"""
Format Converter - Multi-format educational content processor

Converts various input formats (PDF, audio, video, images, presentations)
into processable text for the Pharmacy Professor agent.

Supports:
- PDF: Text extraction with OCR fallback
- Audio/Video: Transcription via Whisper
- Images: OCR + pharmaceutical diagram analysis
- PowerPoint: Slide extraction
- Markdown/Text: Direct processing
"""

import os
import json
import hashlib
import mimetypes
from pathlib import Path
from dataclasses import dataclass, field, asdict
from typing import Optional, List, Dict, Any, Union
from enum import Enum
from datetime import datetime


class InputFormat(Enum):
    """Supported input formats."""
    PDF = "pdf"
    MARKDOWN = "markdown"
    TEXT = "text"
    AUDIO = "audio"
    VIDEO = "video"
    IMAGE = "image"
    POWERPOINT = "powerpoint"
    UNKNOWN = "unknown"


@dataclass
class ConversionResult:
    """Result of format conversion."""
    success: bool
    text: str
    format_detected: InputFormat
    source_file: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    warnings: List[str] = field(default_factory=list)
    processing_time_ms: float = 0.0

    def to_dict(self) -> Dict[str, Any]:
        result = asdict(self)
        result['format_detected'] = self.format_detected.value
        return result


@dataclass
class PageContent:
    """Content from a single page/slide/segment."""
    page_number: int
    text: str
    images: List[str] = field(default_factory=list)
    tables: List[Dict] = field(default_factory=list)


class FormatConverter:
    """
    Multi-format content converter for pharmaceutical education.

    Handles detection and conversion of various content formats into
    processable text while preserving structure and metadata.
    """

    # MIME type to format mapping
    MIME_MAPPINGS = {
        'application/pdf': InputFormat.PDF,
        'text/markdown': InputFormat.MARKDOWN,
        'text/x-markdown': InputFormat.MARKDOWN,
        'text/plain': InputFormat.TEXT,
        'audio/mpeg': InputFormat.AUDIO,
        'audio/wav': InputFormat.AUDIO,
        'audio/mp3': InputFormat.AUDIO,
        'audio/x-wav': InputFormat.AUDIO,
        'video/mp4': InputFormat.VIDEO,
        'video/quicktime': InputFormat.VIDEO,
        'video/x-msvideo': InputFormat.VIDEO,
        'image/png': InputFormat.IMAGE,
        'image/jpeg': InputFormat.IMAGE,
        'image/gif': InputFormat.IMAGE,
        'image/webp': InputFormat.IMAGE,
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': InputFormat.POWERPOINT,
        'application/vnd.ms-powerpoint': InputFormat.POWERPOINT,
    }

    # Extension to format mapping
    EXTENSION_MAPPINGS = {
        '.pdf': InputFormat.PDF,
        '.md': InputFormat.MARKDOWN,
        '.markdown': InputFormat.MARKDOWN,
        '.txt': InputFormat.TEXT,
        '.text': InputFormat.TEXT,
        '.mp3': InputFormat.AUDIO,
        '.wav': InputFormat.AUDIO,
        '.m4a': InputFormat.AUDIO,
        '.ogg': InputFormat.AUDIO,
        '.mp4': InputFormat.VIDEO,
        '.mov': InputFormat.VIDEO,
        '.avi': InputFormat.VIDEO,
        '.mkv': InputFormat.VIDEO,
        '.png': InputFormat.IMAGE,
        '.jpg': InputFormat.IMAGE,
        '.jpeg': InputFormat.IMAGE,
        '.gif': InputFormat.IMAGE,
        '.webp': InputFormat.IMAGE,
        '.pptx': InputFormat.POWERPOINT,
        '.ppt': InputFormat.POWERPOINT,
    }

    def __init__(
        self,
        ocr_enabled: bool = True,
        transcription_model: str = "whisper-1",
        max_file_size_mb: int = 100,
        temp_dir: Optional[Path] = None
    ):
        """
        Initialize format converter.

        Args:
            ocr_enabled: Enable OCR for scanned PDFs and images
            transcription_model: Model for audio/video transcription
            max_file_size_mb: Maximum file size to process
            temp_dir: Directory for temporary files
        """
        self.ocr_enabled = ocr_enabled
        self.transcription_model = transcription_model
        self.max_file_size_mb = max_file_size_mb
        self.temp_dir = temp_dir or Path("/tmp/pharmacy-professor")
        self.temp_dir.mkdir(parents=True, exist_ok=True)

    def detect_format(self, file_path: Union[str, Path]) -> InputFormat:
        """
        Detect input format from file.

        Args:
            file_path: Path to the file

        Returns:
            Detected InputFormat
        """
        path = Path(file_path)

        # Try extension first (most reliable)
        ext = path.suffix.lower()
        if ext in self.EXTENSION_MAPPINGS:
            return self.EXTENSION_MAPPINGS[ext]

        # Try MIME type
        mime_type, _ = mimetypes.guess_type(str(path))
        if mime_type and mime_type in self.MIME_MAPPINGS:
            return self.MIME_MAPPINGS[mime_type]

        # Try reading file header for magic bytes
        if path.exists():
            with open(path, 'rb') as f:
                header = f.read(16)

            # PDF magic bytes
            if header.startswith(b'%PDF'):
                return InputFormat.PDF
            # PNG magic bytes
            if header.startswith(b'\x89PNG'):
                return InputFormat.IMAGE
            # JPEG magic bytes
            if header.startswith(b'\xff\xd8\xff'):
                return InputFormat.IMAGE

        return InputFormat.UNKNOWN

    def convert(
        self,
        file_path: Union[str, Path],
        force_format: Optional[InputFormat] = None
    ) -> ConversionResult:
        """
        Convert file to processable text.

        Args:
            file_path: Path to input file
            force_format: Override format detection

        Returns:
            ConversionResult with extracted text and metadata
        """
        import time
        start_time = time.time()

        path = Path(file_path)
        warnings = []

        # Validate file exists
        if not path.exists():
            return ConversionResult(
                success=False,
                text="",
                format_detected=InputFormat.UNKNOWN,
                source_file=str(path),
                metadata={"error": f"File not found: {path}"}
            )

        # Check file size
        file_size_mb = path.stat().st_size / (1024 * 1024)
        if file_size_mb > self.max_file_size_mb:
            return ConversionResult(
                success=False,
                text="",
                format_detected=InputFormat.UNKNOWN,
                source_file=str(path),
                metadata={"error": f"File too large: {file_size_mb:.1f}MB > {self.max_file_size_mb}MB"}
            )

        # Detect format
        detected_format = force_format or self.detect_format(path)

        # Route to appropriate converter
        converters = {
            InputFormat.PDF: self._convert_pdf,
            InputFormat.MARKDOWN: self._convert_markdown,
            InputFormat.TEXT: self._convert_text,
            InputFormat.AUDIO: self._convert_audio,
            InputFormat.VIDEO: self._convert_video,
            InputFormat.IMAGE: self._convert_image,
            InputFormat.POWERPOINT: self._convert_powerpoint,
        }

        converter = converters.get(detected_format)
        if not converter:
            return ConversionResult(
                success=False,
                text="",
                format_detected=detected_format,
                source_file=str(path),
                metadata={"error": f"Unsupported format: {detected_format.value}"}
            )

        try:
            text, metadata, conv_warnings = converter(path)
            warnings.extend(conv_warnings)

            processing_time = (time.time() - start_time) * 1000

            return ConversionResult(
                success=True,
                text=text,
                format_detected=detected_format,
                source_file=str(path),
                metadata=metadata,
                warnings=warnings,
                processing_time_ms=processing_time
            )

        except Exception as e:
            return ConversionResult(
                success=False,
                text="",
                format_detected=detected_format,
                source_file=str(path),
                metadata={"error": str(e)},
                processing_time_ms=(time.time() - start_time) * 1000
            )

    def _convert_pdf(self, path: Path) -> tuple[str, Dict, List[str]]:
        """Extract text from PDF."""
        warnings = []
        metadata = {
            "type": "pdf",
            "file_name": path.name,
            "file_size_bytes": path.stat().st_size
        }

        text_parts = []

        try:
            # Try PyMuPDF (fitz) first - best for text PDFs
            import fitz

            doc = fitz.open(str(path))
            metadata["page_count"] = len(doc)
            metadata["pdf_title"] = doc.metadata.get("title", "")
            metadata["pdf_author"] = doc.metadata.get("author", "")

            for page_num, page in enumerate(doc, 1):
                page_text = page.get_text()

                if page_text.strip():
                    text_parts.append(f"\n## Page {page_num}\n\n{page_text}")
                elif self.ocr_enabled:
                    # Try OCR for scanned pages
                    ocr_text = self._ocr_page(page)
                    if ocr_text:
                        text_parts.append(f"\n## Page {page_num} (OCR)\n\n{ocr_text}")
                        warnings.append(f"Page {page_num} required OCR")

            doc.close()

        except ImportError:
            # Fallback to pypdf
            try:
                from pypdf import PdfReader

                reader = PdfReader(str(path))
                metadata["page_count"] = len(reader.pages)

                for page_num, page in enumerate(reader.pages, 1):
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        text_parts.append(f"\n## Page {page_num}\n\n{page_text}")

            except ImportError:
                raise ImportError("No PDF library available. Install 'pymupdf' or 'pypdf'")

        return "\n".join(text_parts), metadata, warnings

    def _convert_markdown(self, path: Path) -> tuple[str, Dict, List[str]]:
        """Read markdown file directly."""
        text = path.read_text(encoding='utf-8')
        metadata = {
            "type": "markdown",
            "file_name": path.name,
            "file_size_bytes": path.stat().st_size,
            "line_count": text.count('\n') + 1
        }
        return text, metadata, []

    def _convert_text(self, path: Path) -> tuple[str, Dict, List[str]]:
        """Read text file directly."""
        # Try common encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'ascii']
        text = None
        used_encoding = None

        for encoding in encodings:
            try:
                text = path.read_text(encoding=encoding)
                used_encoding = encoding
                break
            except UnicodeDecodeError:
                continue

        if text is None:
            # Binary fallback
            text = path.read_bytes().decode('utf-8', errors='replace')
            used_encoding = 'utf-8 (with replacements)'

        metadata = {
            "type": "text",
            "file_name": path.name,
            "file_size_bytes": path.stat().st_size,
            "encoding": used_encoding,
            "line_count": text.count('\n') + 1
        }

        warnings = []
        if used_encoding != 'utf-8':
            warnings.append(f"File used {used_encoding} encoding")

        return text, metadata, warnings

    def _convert_audio(self, path: Path) -> tuple[str, Dict, List[str]]:
        """Transcribe audio file."""
        warnings = []
        metadata = {
            "type": "audio",
            "file_name": path.name,
            "file_size_bytes": path.stat().st_size
        }

        # Get audio duration if possible
        try:
            from mutagen import File as AudioFile
            audio = AudioFile(str(path))
            if audio and audio.info:
                metadata["duration_seconds"] = audio.info.length
        except ImportError:
            warnings.append("Install 'mutagen' for audio metadata")
        except Exception:
            pass

        # Transcribe using OpenAI Whisper API or local whisper
        transcript = self._transcribe_audio(path)
        metadata["transcription_model"] = self.transcription_model

        return transcript, metadata, warnings

    def _convert_video(self, path: Path) -> tuple[str, Dict, List[str]]:
        """Extract audio from video and transcribe."""
        warnings = []
        metadata = {
            "type": "video",
            "file_name": path.name,
            "file_size_bytes": path.stat().st_size
        }

        # Extract audio track
        audio_path = self.temp_dir / f"{path.stem}_audio.mp3"

        try:
            import subprocess
            result = subprocess.run([
                'ffmpeg', '-i', str(path),
                '-vn', '-acodec', 'libmp3lame',
                '-y', str(audio_path)
            ], capture_output=True, text=True)

            if result.returncode != 0:
                raise RuntimeError(f"FFmpeg error: {result.stderr}")

            # Transcribe extracted audio
            transcript = self._transcribe_audio(audio_path)

            # Cleanup
            audio_path.unlink(missing_ok=True)

            metadata["transcription_model"] = self.transcription_model
            return transcript, metadata, warnings

        except FileNotFoundError:
            raise ImportError("FFmpeg not found. Install ffmpeg for video processing.")

    def _convert_image(self, path: Path) -> tuple[str, Dict, List[str]]:
        """Extract text from image using OCR."""
        warnings = []
        metadata = {
            "type": "image",
            "file_name": path.name,
            "file_size_bytes": path.stat().st_size
        }

        if not self.ocr_enabled:
            return "[Image - OCR disabled]", metadata, ["OCR disabled"]

        try:
            from PIL import Image
            import pytesseract

            img = Image.open(str(path))
            metadata["image_size"] = f"{img.width}x{img.height}"
            metadata["image_mode"] = img.mode

            # Perform OCR
            text = pytesseract.image_to_string(img)

            # Check if image might be a diagram/chart
            if len(text.strip()) < 50:
                warnings.append("Limited text detected - may be diagram/chart")
                text += "\n\n[Note: This appears to be a diagram or chart with limited text content]"

            return text, metadata, warnings

        except ImportError:
            raise ImportError("Install 'pillow' and 'pytesseract' for image OCR")

    def _convert_powerpoint(self, path: Path) -> tuple[str, Dict, List[str]]:
        """Extract text from PowerPoint presentations."""
        warnings = []
        metadata = {
            "type": "powerpoint",
            "file_name": path.name,
            "file_size_bytes": path.stat().st_size
        }

        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            metadata["slide_count"] = len(prs.slides)

            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                if slide_text:
                    text_parts.append(f"\n## Slide {slide_num}\n\n" + "\n\n".join(slide_text))

            return "\n".join(text_parts), metadata, warnings

        except ImportError:
            raise ImportError("Install 'python-pptx' for PowerPoint processing")

    def _ocr_page(self, page) -> str:
        """Perform OCR on a PDF page."""
        if not self.ocr_enabled:
            return ""

        try:
            from PIL import Image
            import pytesseract
            import io

            # Render page to image
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))

            return pytesseract.image_to_string(img)

        except Exception:
            return ""

    def _transcribe_audio(self, audio_path: Path) -> str:
        """Transcribe audio using available service."""
        # Try OpenAI Whisper API first
        try:
            import openai

            client = openai.OpenAI()
            with open(audio_path, "rb") as audio_file:
                transcript = client.audio.transcriptions.create(
                    model=self.transcription_model,
                    file=audio_file
                )
            return transcript.text

        except ImportError:
            pass
        except Exception:
            pass

        # Fallback to local Whisper
        try:
            import whisper

            model = whisper.load_model("base")
            result = model.transcribe(str(audio_path))
            return result["text"]

        except ImportError:
            raise ImportError("Install 'openai' or 'openai-whisper' for transcription")

    def convert_text_content(self, text: str, source_name: str = "text_input") -> ConversionResult:
        """
        Process raw text content directly.

        Args:
            text: Raw text content
            source_name: Name to identify the source

        Returns:
            ConversionResult
        """
        import time
        start_time = time.time()

        metadata = {
            "type": "text",
            "source": source_name,
            "character_count": len(text),
            "line_count": text.count('\n') + 1,
            "word_count": len(text.split())
        }

        return ConversionResult(
            success=True,
            text=text,
            format_detected=InputFormat.TEXT,
            source_file=source_name,
            metadata=metadata,
            processing_time_ms=(time.time() - start_time) * 1000
        )


def main():
    """CLI interface for format converter."""
    import argparse

    parser = argparse.ArgumentParser(
        description="Convert educational content to processable text"
    )
    parser.add_argument("file", help="File to convert")
    parser.add_argument(
        "--output", "-o",
        help="Output file (default: stdout)"
    )
    parser.add_argument(
        "--format", "-f",
        choices=[f.value for f in InputFormat if f != InputFormat.UNKNOWN],
        help="Force specific format"
    )
    parser.add_argument(
        "--no-ocr",
        action="store_true",
        help="Disable OCR for scanned documents"
    )
    parser.add_argument(
        "--json",
        action="store_true",
        help="Output as JSON with metadata"
    )

    args = parser.parse_args()

    converter = FormatConverter(ocr_enabled=not args.no_ocr)

    force_format = InputFormat(args.format) if args.format else None
    result = converter.convert(args.file, force_format=force_format)

    if args.json:
        output = json.dumps(result.to_dict(), indent=2)
    else:
        if not result.success:
            print(f"Error: {result.metadata.get('error', 'Unknown error')}", file=sys.stderr)
            sys.exit(1)
        output = result.text

    if args.output:
        Path(args.output).write_text(output)
        print(f"Written to {args.output}")
    else:
        print(output)


if __name__ == "__main__":
    import sys
    main()
